"""Legacy text extractors (moved from the API layer).

``extract_text`` dispatches on file suffix: binary office/pdf formats use the
dedicated ``_extract_*`` helpers; everything else is decoded as text — strict
UTF-8 first, then a best-effort charset-detection / GB18030 fallback for
legacy-encoded exports (e.g. GBK CSVs from Chinese Excel).  The remaining
validation (empty / binary-looking / undecodable) is raised as
:class:`LegacyExtractionError` so a bad file fails the parse stage instead of
silently yielding an empty document.
"""

from __future__ import annotations

from io import BytesIO

from lightrag.utils import logger


class LegacyExtractionError(ValueError):
    """Raised when legacy extraction cannot produce usable text."""


def _extract_pdf_pypdf(file_bytes: bytes, password: str | None = None) -> str:
    """Extract PDF content using pypdf (synchronous)."""
    from pypdf import PdfReader  # type: ignore

    pdf_file = BytesIO(file_bytes)
    reader = PdfReader(pdf_file)

    if reader.is_encrypted:
        # Try empty password first (covers permission-only encrypted PDFs)
        decrypt_result = reader.decrypt(password or "")
        if decrypt_result == 0:
            if password:
                raise Exception("Incorrect PDF password")
            else:
                raise Exception("PDF is encrypted but no password provided")

    content = ""
    for page in reader.pages:
        content += page.extract_text() + "\n"
    return content


def _extract_docx(file_bytes: bytes) -> str:
    """Extract DOCX content including tables in document order (synchronous)."""
    from docx import Document  # type: ignore
    from docx.table import Table  # type: ignore
    from docx.text.paragraph import Paragraph  # type: ignore

    docx_file = BytesIO(file_bytes)
    doc = Document(docx_file)

    def escape_cell(cell_value: str | None) -> str:
        if cell_value is None:
            return ""
        text = str(cell_value)
        return (
            text.replace("\\", "\\\\")
            .replace("\t", "&emsp;&emsp;")
            .replace("\r\n", "<br>")
            .replace("\r", "<br>")
            .replace("\n", "<br>")
        )

    content_parts = []
    in_table = False
    for element in doc.element.body:
        if element.tag.endswith("p"):
            if in_table:
                content_parts.append("")
                in_table = False
            paragraph = Paragraph(element, doc)
            content_parts.append(paragraph.text)
        elif element.tag.endswith("tbl"):
            if content_parts and not in_table:
                content_parts.append("")
            in_table = True
            table = Table(element, doc)
            for row in table.rows:
                row_text = [escape_cell(cell.text) for cell in row.cells]
                if any(cell for cell in row_text):
                    content_parts.append("\t".join(row_text))
    return "\n".join(content_parts)


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract PPTX content (synchronous)."""
    from pptx import Presentation  # type: ignore

    pptx_file = BytesIO(file_bytes)
    prs = Presentation(pptx_file)
    content = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content


def _extract_xlsx(file_bytes: bytes) -> str:
    """Extract XLSX content in tab-delimited format with sheet separators."""
    from openpyxl import load_workbook  # type: ignore

    xlsx_file = BytesIO(file_bytes)
    wb = load_workbook(xlsx_file)

    def escape_cell(cell_value: str | int | float | None) -> str:
        if cell_value is None:
            return ""
        text = str(cell_value)
        return (
            text.replace("\\", "\\\\")
            .replace("\t", "\\t")
            .replace("\r\n", "\\n")
            .replace("\r", "\\n")
            .replace("\n", "\\n")
        )

    def escape_sheet_title(title: str) -> str:
        return str(title).replace("\n", " ").replace("\t", " ").replace("\r", " ")

    content_parts: list[str] = []
    sheet_separator = "=" * 20

    for idx, sheet in enumerate(wb):
        if idx > 0:
            content_parts.append("")
        safe_title = escape_sheet_title(sheet.title)
        content_parts.append(f"{sheet_separator} Sheet: {safe_title} {sheet_separator}")
        max_columns = sheet.max_column if sheet.max_column else 0
        for row in sheet.iter_rows(values_only=True):
            row_parts = []
            for col_idx in range(max_columns):
                if col_idx < len(row):
                    row_parts.append(escape_cell(row[col_idx]))
                else:
                    row_parts.append("")
            if all(part == "" for part in row_parts):
                content_parts.append("")
            else:
                content_parts.append("\t".join(row_parts))
    content_parts.append(sheet_separator)
    return "\n".join(content_parts)


# Suffixes (without dot) routed to dedicated binary extractors.
_BINARY_EXTRACTORS = {
    "pdf": _extract_pdf_pypdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "xlsx": _extract_xlsx,
}


def _decode_text(file_bytes: bytes) -> str:
    """Decode text bytes with the upload-path validation, raised on failure.

    Strict UTF-8 (BOM stripped via ``utf-8-sig``) is tried first; non-UTF-8
    input falls back to :func:`_decode_non_utf8_fallback` instead of failing
    outright.
    """
    try:
        content = file_bytes.decode("utf-8-sig")
    except UnicodeDecodeError as utf8_error:
        content = _decode_non_utf8_fallback(file_bytes, utf8_error)
    if not content or len(content.strip()) == 0:
        raise LegacyExtractionError("File contains no content or only whitespace")
    if content.startswith("b'") or content.startswith('b"'):
        raise LegacyExtractionError(
            "File appears to contain binary data representation instead of text"
        )
    return content


def _decode_non_utf8_fallback(file_bytes: bytes, utf8_error: UnicodeDecodeError) -> str:
    """Best-effort decode for legacy-encoded text (e.g. GBK CSV exports).

    Tries strict GB18030 first: this deployment's non-UTF-8 uploads are
    overwhelmingly GBK/GB2312 exports from Chinese Excel, and statistical
    charset detection demonstrably misdetects short CJK snippets (observed:
    a GBK CSV decoded as Korean cp949, yielding garbage).  GB18030 is a
    superset of GBK/GB2312 and rejects (raises on) byte sequences that are
    not valid Chinese text, in which case we fall through to
    charset-normalizer for other legacy encodings.
    """
    try:
        content = file_bytes.decode("gb18030")
    except UnicodeDecodeError:
        pass
    else:
        logger.warning(
            "Non-UTF-8 text decoded with GB18030 fallback; "
            "consider converting the file to UTF-8"
        )
        return content

    try:
        from charset_normalizer import from_bytes
    except ImportError:
        from_bytes = None

    if from_bytes is not None:
        try:
            best = from_bytes(file_bytes).best()
        except Exception:  # detection must never break ingestion
            best = None
        if best is not None and best.encoding:
            logger.info(
                "Non-UTF-8 text decoded via charset detection "
                f"(encoding={best.encoding})"
            )
            return str(best)

    raise LegacyExtractionError(
        "File is not valid UTF-8 encoded text and could not be decoded "
        "with any fallback encoding. Please convert it to UTF-8 before "
        f"processing: {utf8_error}"
    ) from utf8_error


def extract_text(
    file_bytes: bytes, suffix: str, *, pdf_password: str | None = None
) -> str:
    """Extract plain text from ``file_bytes`` based on ``suffix`` (no dot).

    Synchronous; callers run it in a thread.  Raises
    :class:`LegacyExtractionError` (or the extractor's own exception) on
    failure.
    """
    suffix = suffix.lower().lstrip(".")
    extractor = _BINARY_EXTRACTORS.get(suffix)
    if extractor is _extract_pdf_pypdf:
        return _extract_pdf_pypdf(file_bytes, pdf_password)
    if extractor is not None:
        return extractor(file_bytes)
    return _decode_text(file_bytes)
